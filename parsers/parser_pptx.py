from pptx import Presentation
from parsers.api import TranslateAPI
from parsers.db import DB


class ParserPPTX(object):
    def __init__(self, row_id, src_file, des_file, src_lang, des_lang):
        self.row_id = row_id
        self.src_file = src_file
        self.des_file = des_file
        self.src_lang = src_lang
        self.des_lang = des_lang

    def calculate_total_progress(self):
        total = 0
        prs = Presentation(self.src_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        total = total + 1
        return total

    def parse(self):
        total = self.calculate_total_progress()
        current = 0
        percent = 0
        prs = Presentation(self.src_file)
        trans = TranslateAPI()
        db = DB()
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = trans.translate(self.src_lang, self.des_lang, run.text)
                        current = current + 1
                        if percent != int(current * 100 / total):
                            percent = int(current * 100 / total)
                            db.update_record_progress(self.row_id, percent)
        db.update_record_progress(self.row_id, 100)
        prs.save(self.des_file)
        trans.close()
        db.close()
